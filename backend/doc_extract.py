"""
事前学習資料（PDF/PPT/Word）からの氏名候補抽出。

設計方針（docs/app-design.md 5章）:
- テキスト層があるページ／スライドのみ対象（OCRはスコープ外）
- 表構造に「氏名」「ふりがな」らしき列があれば、それを最優先で抽出
  （表彰リスト等はこの形式が多く、読みも一緒に取れて確実なため）
- 自由記述の文章からは日本語NER（GiNZA）で人名候補を検出する
- 読みが取れない場合は pykakasi で推定するが、あくまで「候補」であり、
  最終確認は主催者が行う前提（人名の読みは推定が外れやすいため）
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pykakasi import kakasi

_kks = kakasi()

NAME_HEADER_HINTS = ["氏名", "名前", "お名前", "受賞者", "登壇者"]
READING_HEADER_HINTS = ["ふりがな", "フリガナ", "よみがな", "読み", "かな"]

# GiNZAは重いモデルのため、実際に使われるまで読み込みを遅らせる
_ginza_nlp = None
_ginza_load_failed = False


def _get_ginza():
    global _ginza_nlp, _ginza_load_failed
    if _ginza_nlp is not None or _ginza_load_failed:
        return _ginza_nlp
    try:
        import spacy
        _ginza_nlp = spacy.load("ja_ginza")
    except Exception:
        _ginza_load_failed = True
        _ginza_nlp = None
    return _ginza_nlp


def guess_reading(text: str) -> str:
    """漢字の氏名から読み（ひらがな）を推定する。あくまで候補であり、
    主催者による確認・修正が前提。
    """
    return "".join(item["hira"] for item in _kks.convert(text))


def _find_name_reading_columns(header_row: list[str]) -> Optional[tuple[int, Optional[int]]]:
    """表のヘッダー行から「氏名列」「ふりがな列」のインデックスを探す"""
    name_idx = None
    reading_idx = None
    for i, cell in enumerate(header_row):
        cell = (cell or "").strip()
        if name_idx is None and any(h in cell for h in NAME_HEADER_HINTS):
            name_idx = i
        if reading_idx is None and any(h in cell for h in READING_HEADER_HINTS):
            reading_idx = i
    if name_idx is None:
        return None
    return name_idx, reading_idx


def _candidates_from_table(rows: list[list[str]]) -> list[dict]:
    if not rows or len(rows) < 2:
        return []
    cols = _find_name_reading_columns(rows[0])
    if cols is None:
        return []
    name_idx, reading_idx = cols

    candidates = []
    for row in rows[1:]:
        if name_idx >= len(row):
            continue
        name = (row[name_idx] or "").strip()
        if not name:
            continue
        reading = ""
        if reading_idx is not None and reading_idx < len(row):
            reading = (row[reading_idx] or "").strip()
        candidates.append({
            "name": name,
            "reading": reading or guess_reading(name),
            "readingGuessed": not bool(reading),
            "source": "table",
        })
    return candidates


def _candidates_from_text(text: str) -> list[dict]:
    nlp = _get_ginza()
    if nlp is None or not text.strip():
        return []
    doc = nlp(text)
    seen = set()
    candidates = []
    for ent in doc.ents:
        if ent.label_ != "Person":
            continue
        name = ent.text.strip()
        if not name or name in seen:
            continue
        seen.add(name)
        candidates.append({
            "name": name,
            "reading": guess_reading(name),
            "readingGuessed": True,
            "source": "ner",
        })
    return candidates


def _extract_pdf(path: str) -> tuple[list[dict], int]:
    import pdfplumber

    candidates: list[dict] = []
    unreadable_pages = 0
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if not text.strip():
                unreadable_pages += 1
                continue
            for table in page.extract_tables() or []:
                candidates.extend(_candidates_from_table(table))
            candidates.extend(_candidates_from_text(text))
    return candidates, unreadable_pages


def _extract_docx(path: str) -> tuple[list[dict], int]:
    import docx

    document = docx.Document(path)
    candidates: list[dict] = []

    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        candidates.extend(_candidates_from_table(rows))

    free_text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    candidates.extend(_candidates_from_text(free_text))
    return candidates, 0


def _extract_pptx(path: str) -> tuple[list[dict], int]:
    from pptx import Presentation

    prs = Presentation(path)
    candidates: list[dict] = []

    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                candidates.extend(_candidates_from_table(rows))
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t.strip():
                    texts.append(t)
        candidates.extend(_candidates_from_text("\n".join(texts)))
    return candidates, 0


def extract_name_candidates(path: str, filename: str) -> dict:
    """資料ファイルから氏名候補を抽出する。

    戻り値: {"candidates": [...], "unreadablePages": N, "engine": "ginza"|"table-only"}
    """
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        raw, unreadable = _extract_pdf(path)
    elif ext == ".docx":
        raw, unreadable = _extract_docx(path)
    elif ext == ".pptx":
        raw, unreadable = _extract_pptx(path)
    else:
        raise ValueError(f"unsupported file type: {ext}")

    # 同じ名前が複数箇所から出た場合、表由来（読みが確実）を優先して重複排除
    dedup: dict[str, dict] = {}
    for c in raw:
        existing = dedup.get(c["name"])
        if existing is None or (existing["source"] == "ner" and c["source"] == "table"):
            dedup[c["name"]] = c

    return {
        "candidates": list(dedup.values()),
        "unreadablePages": unreadable,
        "nerAvailable": _get_ginza() is not None,
    }
